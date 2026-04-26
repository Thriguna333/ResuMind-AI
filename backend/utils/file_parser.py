import re
import fitz
import spacy
from docx import Document
from pptx import Presentation

nlp = spacy.load("en_core_web_sm")


def clean_text(text):
    lines = text.split("\n")
    cleaned = []

    for line in lines:
        line = line.strip()

        # remove useless lines
        if not line:
            continue
        if line.isdigit():
            continue
        if len(line) < 3:
            continue

        cleaned.append(line)

    return "\n".join(cleaned)


def extract_key_value_pairs(text):
    kv_pairs = {}

    lines = text.split("\n")

    allowed_keys = [
        "phone", "email", "linkedin", "github",
        "cgpa", "expected graduation", "location"
    ]

    for line in lines:
        line = line.strip().lower()

        # Case 1: Key: Value
        match = re.match(r"^([a-zA-Z ]+)\s*[:\-]\s*(.+)$", line)

        if match:
            key = match.group(1).strip()
            value = match.group(2).strip()

            if key in allowed_keys:
                kv_pairs[key.title()] = value

            continue

        # Case 2: Key Value
        match2 = re.match(
            r"^(phone|email|linkedin|github|cgpa|expected graduation|location)\s+(.+)$",
            line
        )

        if match2:
            key = match2.group(1).strip()
            value = match2.group(2).strip()

            kv_pairs[key.title()] = value

    return kv_pairs


def extract_entities(text):
    doc = nlp(text[:100000])
    entities = []

    for ent in doc.ents:
        entities.append({
            "text": ent.text,
            "label": ent.label_,
            "start_char": ent.start_char,
            "end_char": ent.end_char
        })

    return entities


def extract_pdf_headings(page):
    headings = []
    blocks = page.get_text("dict")["blocks"]

    for block in blocks:
        if "lines" not in block:
            continue

        for line in block["lines"]:
            line_text = ""
            max_font_size = 0
            is_bold = False

            for span in line["spans"]:
                text = span["text"].strip()
                if not text:
                    continue

                line_text += text + " "
                max_font_size = max(max_font_size, span["size"])

                font_name = span.get("font", "").lower()
                if "bold" in font_name:
                    is_bold = True

            line_text = line_text.strip()

            if not line_text:
                continue

            if max_font_size >= 14 or (is_bold and len(line_text.split()) <= 12):
                headings.append({
                    "heading": line_text,
                    "font_size": round(max_font_size, 2),
                    "is_bold": is_bold
                })

    return headings


def extract_pdf_tables(page):
    tables = []

    try:
        found_tables = page.find_tables()

        for table_index, table in enumerate(found_tables.tables, start=1):
            data = table.extract()

            tables.append({
                "table_number": table_index,
                "rows": data,
                "row_count": len(data),
                "column_count": len(data[0]) if data else 0
            })

    except Exception:
        pass

    return tables


def get_section(text, start_heading, end_headings):
    lines = text.split("\n")
    capture = False
    section_lines = []

    for line in lines:
        clean = line.strip()

        if clean.upper() == start_heading.upper():
            capture = True
            continue

        if capture and clean.upper() in [h.upper() for h in end_headings]:
            break

        if capture and clean:
            section_lines.append(clean)

    return section_lines


def extract_resume_intelligence(text, key_value_pairs):
    lines = [line.strip() for line in text.split("\n") if line.strip()]

    name = lines[0].title() if lines else ""
    location = lines[1] if len(lines) > 1 else ""

    skills_lines = get_section(
        text,
        "SKILLS",
        ["PROJECTS", "EDUCATION", "CERTIFICATIONS & ACTIVITIES", "LANGUAGES"]
    )

    skills = []
    for line in skills_lines:
        if ":" in line:
            _, value = line.split(":", 1)
            items = [item.strip() for item in value.split(",")]
            skills.extend(items)

    projects_lines = get_section(
        text,
        "PROJECTS",
        ["EDUCATION", "CERTIFICATIONS & ACTIVITIES", "LANGUAGES"]
    )

    projects = []
    for line in projects_lines:
        if not line.startswith("•") and "–" in line:
            projects.append(line)

    education_lines = get_section(
        text,
        "EDUCATION",
        ["CERTIFICATIONS & ACTIVITIES", "LANGUAGES"]
    )

    education = {
        "degree": education_lines[0] if len(education_lines) > 0 else "",
        "college": education_lines[1] if len(education_lines) > 1 else "",
        "cgpa": key_value_pairs.get("Cgpa", key_value_pairs.get("CGPA", "")),
        "expected_graduation": key_value_pairs.get("Expected Graduation", "")
    }

    return {
        "name": name,
        "location": location,
        "phone": key_value_pairs.get("Phone", ""),
        "email": key_value_pairs.get("Email", ""),
        "linkedin": key_value_pairs.get("Linkedin", key_value_pairs.get("LinkedIn", "")),
        "skills": skills,
        "education": education,
        "projects": projects
    }



def parse_pdf(file_path):
    pages = []
    all_headings = []
    all_tables = []
    all_key_value_pairs = {}
    all_entities = []

    doc = fitz.open(file_path)

    for page_number, page in enumerate(doc, start=1):
        text = page.get_text()
        headings = extract_pdf_headings(page)
        tables = extract_pdf_tables(page)
        key_value_pairs = extract_key_value_pairs(text)
        entities = extract_entities(text)

        for heading in headings:
            heading["page_number"] = page_number
            all_headings.append(heading)

        for table in tables:
            table["page_number"] = page_number
            all_tables.append(table)

        for entity in entities:
            entity["page_number"] = page_number
            all_entities.append(entity)

        all_key_value_pairs.update(key_value_pairs)

        pages.append({
            "page_number": page_number,
            "text": text,
            "word_count": len(text.split()),
            "char_count": len(text),
            "headings": headings,
            "tables": tables,
            "key_value_pairs": key_value_pairs,
            "entities": entities
        })

    full_text = "\n".join([page["text"] for page in pages])
    full_text = clean_text(full_text)

    return {
        "full_text": full_text,
        "pages": pages,
        "headings": all_headings,
        "tables": all_tables,
        "key_value_pairs": all_key_value_pairs,
        "entities": all_entities,
        "resume_intelligence": extract_resume_intelligence(full_text, all_key_value_pairs)
    }


def parse_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    text = clean_text(text)

    headings = []
    tables = []
    key_value_pairs = extract_key_value_pairs(text)
    entities = extract_entities(text)

    for para in doc.paragraphs:
        if para.style and para.style.name.startswith("Heading"):
            headings.append({
                "heading": para.text,
                "style": para.style.name,
                "page_number": 1
            })

    for table_index, table in enumerate(doc.tables, start=1):
        rows = []

        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])

        tables.append({
            "table_number": table_index,
            "page_number": 1,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(rows[0]) if rows else 0
        })

    return {
        "full_text": text,
        "pages": [
            {
                "page_number": 1,
                "text": text,
                "word_count": len(text.split()),
                "char_count": len(text),
                "headings": headings,
                "tables": tables,
                "key_value_pairs": key_value_pairs,
                "entities": entities,
                "resume_intelligence": extract_resume_intelligence(text, key_value_pairs)
            }
        ],
        "headings": headings,
        "tables": tables,
        "key_value_pairs": key_value_pairs,
        "entities": entities,
        "resume_intelligence": extract_resume_intelligence(text, key_value_pairs)
    }


def parse_pptx(file_path):
    prs = Presentation(file_path)
    pages = []
    all_headings = []
    all_tables = []
    all_key_value_pairs = {}
    all_entities = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        slide_headings = []
        slide_tables = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                slide_text += text + "\n"

                if shape.is_placeholder and text:
                    slide_headings.append({
                        "heading": text,
                        "page_number": slide_number
                    })

            if shape.has_table:
                rows = []

                for row in shape.table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])

                slide_tables.append({
                    "table_number": len(slide_tables) + 1,
                    "page_number": slide_number,
                    "rows": rows,
                    "row_count": len(rows),
                    "column_count": len(rows[0]) if rows else 0
                })

        key_value_pairs = extract_key_value_pairs(slide_text)
        entities = extract_entities(slide_text)

        all_key_value_pairs.update(key_value_pairs)
        all_headings.extend(slide_headings)
        all_tables.extend(slide_tables)

        for entity in entities:
            entity["page_number"] = slide_number
            all_entities.append(entity)

        pages.append({
            "page_number": slide_number,
            "text": slide_text,
            "word_count": len(slide_text.split()),
            "char_count": len(slide_text),
            "headings": slide_headings,
            "tables": slide_tables,
            "key_value_pairs": key_value_pairs,
            "entities": entities
        })

    full_text = "\n".join([page["text"] for page in pages])
    full_text = clean_text(full_text)

    return {
        "full_text": full_text,
        "pages": pages,
        "headings": all_headings,
        "tables": all_tables,
        "key_value_pairs": all_key_value_pairs,
        "entities": all_entities,
        "resume_intelligence": extract_resume_intelligence(full_text, all_key_value_pairs)
    }


def parse_file(file_path, file_type):
    if file_type == "pdf":
        return parse_pdf(file_path)
    elif file_type == "docx":
        return parse_docx(file_path)
    elif file_type == "pptx":
        return parse_pptx(file_path)
    else:
        return {
            "full_text": "",
            "pages": [],
            "headings": [],
            "tables": [],
            "key_value_pairs": {},
            "entities": [],
            "resume_intelligence": {},
            "error": "Unsupported file type"
        }
    

def analyze_resume(text):
    score = 0
    feedback = []

    text_lower = text.lower()

    if "python" in text_lower:
        score += 20
    else:
        feedback.append("Add Python skill")

    if "machine learning" in text_lower:
        score += 20
    else:
        feedback.append("Add Machine Learning projects")

    if "internship" in text_lower:
        score += 20
    else:
        feedback.append("Add internship experience")

    if "project" in text_lower or "projects" in text_lower:
        score += 20
    else:
        feedback.append("Add more projects")

    if "sql" in text_lower:
        score += 20
    else:
        feedback.append("Add SQL skills")

    return {
        "score": score,
        "feedback": feedback
    }